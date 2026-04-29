from __future__ import annotations

import hashlib
import re
from pathlib import Path

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

from src.models.schemas import DocumentMetadata, ParsedDocument


class MinerUParser:
    """MinerU-first parser with safe format fallbacks."""

    def __init__(self, mineru_enabled: bool = True):
        self.mineru_enabled = mineru_enabled

    def parse(self, source_path: str | Path) -> ParsedDocument:
        path = Path(source_path)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {path}")
        doc_id = hashlib.sha1(path.read_bytes()).hexdigest()[:16]
        ext = path.suffix.lower()
        metadata = DocumentMetadata(doc_id=doc_id, source_path=path, source_type=ext.lstrip("."))

        if self.mineru_enabled and ext == ".pdf":
            mineru_output = self._parse_pdf_with_mineru(path)
            if mineru_output:
                markdown = self._to_markdown_from_text(mineru_output)
                return ParsedDocument(metadata=metadata, markdown=markdown, plain_text=mineru_output)

        if ext == ".pdf":
            text = self._parse_pdf_fallback(path)
        elif ext == ".docx":
            text = self._parse_docx(path)
        elif ext == ".pptx":
            text = self._parse_pptx(path)
        elif ext in {".md", ".txt"}:
            text = path.read_text(encoding="utf-8", errors="ignore")
        else:
            text = path.read_text(encoding="utf-8", errors="ignore")

        markdown = self._to_markdown_from_text(text)
        return ParsedDocument(metadata=metadata, markdown=markdown, plain_text=text)

    def _parse_pdf_with_mineru(self, path: Path) -> str:
        """Try MinerU API if installed; fallback silently when unavailable."""
        try:
            # MinerU commonly exposes these modules (renamed from magic-pdf).
            from magic_pdf.data.data_reader_writer import FileBasedDataWriter
            from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze
            from magic_pdf.pipe.UNIPipe import UNIPipe
        except Exception:
            return ""

        try:
            output_dir = path.parent / ".mineru_tmp"
            output_dir.mkdir(parents=True, exist_ok=True)
            writer = FileBasedDataWriter(str(output_dir))
            infer_result = doc_analyze(path.read_bytes(), ocr=True)
            pipe = UNIPipe(path.read_bytes(), infer_result, image_writer=writer)
            pipe.pipe_parse()
            md = pipe.pipe_mk_markdown("")
            return md or ""
        except Exception:
            return ""

    def _parse_pdf_fallback(self, path: Path) -> str:
        reader = PdfReader(str(path))
        pages = []
        for idx, page in enumerate(reader.pages):
            content = page.extract_text() or ""
            pages.append(f"# Page {idx + 1}\n{content.strip()}")
        return "\n\n".join(pages)

    def _parse_docx(self, path: Path) -> str:
        doc = DocxDocument(str(path))
        lines: list[str] = []
        for p in doc.paragraphs:
            txt = p.text.strip()
            if txt:
                lines.append(txt)
        return "\n".join(lines)

    def _parse_pptx(self, path: Path) -> str:
        prs = Presentation(str(path))
        out: list[str] = []
        for idx, slide in enumerate(prs.slides):
            out.append(f"# Slide {idx + 1}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text.strip())
        return "\n".join(line for line in out if line)

    def _to_markdown_from_text(self, text: str) -> str:
        lines = [ln.rstrip() for ln in text.splitlines()]
        normalized: list[str] = []
        for line in lines:
            if not line:
                normalized.append("")
                continue
            if re.match(r"^\s*(page|slide)\s+\d+", line, re.IGNORECASE):
                normalized.append(f"## {line.strip()}")
            else:
                normalized.append(line.strip())
        return "\n".join(normalized).strip() + "\n"
